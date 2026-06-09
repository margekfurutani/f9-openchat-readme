# -*- coding: utf-8 -*-
"""OpenChat IVS ドキュメント3本を PowerPoint 化するスクリプト。
   観点: 「この場面ではこのSDK/APIを使う」"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- テーマ色 ----
NAVY   = RGBColor(0x1F, 0x33, 0x55)
BLUE   = RGBColor(0x2E, 0x6F, 0xB5)
ACCENT = RGBColor(0xE8, 0x7A, 0x1E)
LIGHT  = RGBColor(0xEF, 0xF3, 0xF8)
GRAY   = RGBColor(0x55, 0x5F, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x27, 0x2E)
FONT   = "Meiryo"

SW, SH = Inches(13.333), Inches(7.5)  # 16:9


def _set_font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def add_title_slide(prs, kicker, title, subtitle, tag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    bg.shadow.inherit = False
    # 左アクセント帯
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.25), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False
    # kicker
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11), Inches(0.6)).text_frame
    r = tb.paragraphs[0].add_run(); r.text = kicker; _set_font(r, 18, ACCENT, True)
    # title
    tb = s.shapes.add_textbox(Inches(0.85), Inches(2.7), Inches(11.6), Inches(1.8)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = title; _set_font(r, 40, WHITE, True)
    # subtitle
    tb = s.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11), Inches(1.2)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = subtitle; _set_font(r, 18, RGBColor(0xC7,0xD3,0xE2))
    # tag chip
    chip = s.shapes.add_shape(1, Inches(0.9), Inches(6.2), Inches(4.2), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = BLUE; chip.line.fill.background()
    chip.shadow.inherit = False
    ctf = chip.text_frame; ctf.word_wrap = True
    ctf.margin_top = Pt(2); ctf.margin_bottom = Pt(2)
    r = ctf.paragraphs[0].add_run(); r.text = tag; _set_font(r, 13, WHITE, True)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return s


def _header(s, title, idx=None):
    band = s.shapes.add_shape(1, 0, 0, SW, Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
    band.shadow.inherit = False
    bar = s.shapes.add_shape(1, 0, Inches(1.0), SW, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(11.8), Inches(0.7)).text_frame
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tb.paragraphs[0].add_run(); r.text = title; _set_font(r, 26, WHITE, True)
    if idx:
        nb = s.shapes.add_textbox(Inches(12.2), Inches(0.18), Inches(1.0), Inches(0.7)).text_frame
        nb.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = nb.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = idx; _set_font(r, 12, RGBColor(0x9F,0xB0,0xC4))


def add_bullets_slide(prs, title, bullets, idx=None):
    """bullets: list of (level, text, kind) ; kind: 'h'(見出し) / 't'(通常) / 'c'(コード/強調)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    _header(s, title, idx)
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.9)).text_frame
    box.word_wrap = True
    first = True
    for level, text, kind in bullets:
        p = box.paragraphs[0] if first else box.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(6); p.space_before = Pt(2)
        marker = ""
        if kind == 'h':
            p.space_before = Pt(10)
            r = p.add_run(); r.text = text; _set_font(r, 19, BLUE, True)
        elif kind == 'c':
            marker = "▶ " if level == 0 else "• "
            r = p.add_run(); r.text = marker + text
            _set_font(r, 15, DARK, False)
            r.font.name = "Consolas"
        else:
            marker = ("• " if level == 0 else "– ")
            r = p.add_run(); r.text = marker + text
            _set_font(r, 16 if level == 0 else 15, DARK if level == 0 else GRAY)
    return s


def add_table_slide(prs, title, headers, rows, idx=None, col_w=None, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    _header(s, title, idx)
    nrows, ncols = len(rows) + 1, len(headers)
    top = Inches(1.35); left = Inches(0.6)
    width = Inches(12.13); height = Inches(0.5 * nrows)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Emu(int(int(width) * w / total))
    # header
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = h; _set_font(r, 13, WHITE, True)
    # body
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            mono = val and (val[0] in "`_" or "(" in val and ")" in val and val.islower())
            _set_font(r, 12, DARK, c == 0)
    if note:
        nb = s.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.5)).text_frame
        nb.word_wrap = True
        r = nb.paragraphs[0].add_run(); r.text = "※ " + note; _set_font(r, 12, GRAY)
    return s


def add_flow_slide(prs, title, steps, idx=None):
    """steps: list of (label, detail, highlight_bool)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    _header(s, title, idx)
    n = len(steps)
    top = Inches(1.5); gap = Inches(0.12)
    avail = Inches(5.6)
    h = Emu(int((int(avail) - int(gap) * (n - 1)) / n))
    y = int(top)
    for label, detail, hl in steps:
        box = s.shapes.add_shape(1, Inches(0.8), Emu(y), Inches(11.7), h)
        box.fill.solid(); box.fill.fore_color.rgb = (ACCENT if hl else LIGHT)
        box.line.color.rgb = (ACCENT if hl else RGBColor(0xCB,0xD6,0xE2)); box.line.width = Pt(1)
        box.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(14)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = label
        _set_font(r, 16, WHITE if hl else NAVY, True)
        if detail:
            r2 = p.add_run(); r2.text = "   " + detail
            _set_font(r2, 13, WHITE if hl else GRAY, False)
        y += int(h) + int(gap)
    return s


CODE_BG   = RGBColor(0x1E, 0x21, 0x2B)
CODE_TX   = RGBColor(0xE6, 0xE9, 0xEF)
CODE_CMT  = RGBColor(0x7E, 0xC6, 0x99)   # コメント=緑
CODE_KEY  = RGBColor(0x6A, 0xB0, 0xF3)   # 強調行=青
CODE_NUM  = RGBColor(0x6B, 0x74, 0x84)   # 行番号


def _is_comment(line):
    t = line.lstrip()
    return t.startswith("//") or t.startswith("#") or t.startswith("*") \
        or t.startswith("/*") or t.startswith("/**")


def add_code_slide(prs, title, intro, code_lines, idx=None, caption=None, hl=None):
    """intro: list[str] 短い説明 / code_lines: list[str] / hl: set[int] 強調行(0始まり)"""
    hl = hl or set()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    _header(s, title, idx)
    top = 1.3
    if intro:
        ib = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(0.9)).text_frame
        ib.word_wrap = True
        first = True
        for t in intro:
            p = ib.paragraphs[0] if first else ib.add_paragraph()
            first = False
            r = p.add_run(); r.text = "• " + t; _set_font(r, 14, GRAY)
            p.space_after = Pt(2)
        top += 0.35 * len(intro) + 0.25
    # コードボックス
    box_h = 7.15 - top - (0.45 if caption else 0.1)
    panel = s.shapes.add_shape(2, Inches(0.7), Inches(top), Inches(11.93), Inches(box_h))
    panel.fill.solid(); panel.fill.fore_color.rgb = CODE_BG
    panel.line.color.rgb = RGBColor(0x3A,0x40,0x4E); panel.line.width = Pt(1)
    panel.shadow.inherit = False
    tf = panel.text_frame; tf.word_wrap = False
    tf.margin_left = Pt(14); tf.margin_right = Pt(8)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    # 行数からフォント自動調整
    n = len(code_lines)
    fs = 13 if n <= 16 else (12 if n <= 20 else (11 if n <= 24 else 10))
    first = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(0); p.space_before = Pt(0); p.line_spacing = 1.0
        color = CODE_CMT if _is_comment(line) else (CODE_KEY if i in hl else CODE_TX)
        r = p.add_run(); r.text = line if line else " "
        r.font.size = Pt(fs); r.font.name = "Consolas"
        r.font.color.rgb = color
        r.font.bold = (i in hl)
    if caption:
        cb = s.shapes.add_textbox(Inches(0.7), Inches(top + box_h + 0.05), Inches(11.93), Inches(0.4)).text_frame
        cb.word_wrap = True
        r = cb.paragraphs[0].add_run(); r.text = caption; _set_font(r, 12, ACCENT, True)
    return s


def save(prs, path):
    prs.slide_width = SW; prs.slide_height = SH
    prs.save(path)
    print("wrote", path)


def base():
    p = Presentation(); p.slide_width = SW; p.slide_height = SH; return p


# ============================================================
# DECK 1: Flax サーバ
# ============================================================
def deck_server():
    p = base()
    add_title_slide(p,
        "OpenChat × AWS IVS / サーバ編",
        "Flaxサーバ ─ IVS処理順",
        "Stageの作成・トークン発行・削除をどの場面でどう呼ぶか",
        "AWS SDK for Java  ivsrealtime  /  region: ap-northeast-1")

    add_bullets_slide(p, "機能概要", [
        (0, "1人のパフォーマが配信し、複数の会員・ゲストが視聴する一方向ライブ配信", 't'),
        (0, "映像配信は AWS IVS リアルタイムストリーミングの「Stage」を使用", 't'),
        (0, "Flaxの責務", 'h'),
        (0, "Stageの作成・トークン発行・削除を SDK 経由で管理", 't'),
        (0, "配信状態をメモリ上（シングルトン）で保持", 't'),
        (0, "主要コンポーネント", 'h'),
        (0, "IvsRealtimeUtils ─ IVS SDK呼び出しのラッパー", 'c'),
        (0, "OpenChatRoomManager ─ 配信のシングルトン管理", 'c'),
        (0, "OpenChatRoom ─ 1配信の状態（performer/stageArn/視聴者）", 'c'),
        (0, "Performer/Viewer WsClient ─ 配信用・視聴用WSエンドポイント", 'c'),
    ], "01 server / 2")

    add_table_slide(p, "この場面でこのIVS API ─ IvsRealtimeUtils",
        ["メソッド", "IVS API", "用途", "トークン有効期限"],
        [
            ["createStage(name)", "CreateStage", "配信用Stage作成＋配信トークン取得", "publish 1分"],
            ["createSubscribeToken(arn)", "CreateParticipantToken", "視聴トークン発行", "subscribe 1分"],
            ["deleteStage(arn)", "DeleteStage", "Stage削除", "-"],
            ["listStages()", "ListStages", "孤児Stage検出用の全件取得", "-"],
        ],
        "01 server / 3", col_w=[3, 2.4, 3.6, 1.8],
        note="トークン1分は漏洩時の不正利用防止。join時のみ使用し、継続には不要。Stage名は openchat-{performerCode}-{timestamp}")

    add_flow_slide(p, "IVS処理順サマリ ─ IVSを呼ぶのは4箇所だけ", [
        ("① 配信開始", "performer WS接続 → createStage() → publishToken返却", True),
        ("② 視聴開始", "viewer WS接続(配信中のみ) → createSubscribeToken() → subscribeToken返却", True),
        ("③ 視聴終了", "viewer切断 → removeViewer()  ※IVS呼び出しなし", False),
        ("④ 配信終了", "performer切断 → 全視聴者へEND通知 → deleteStage()", True),
        ("⑤ 孤児削除", "cron毎時 → listStages() → 管理外Stageを deleteStage()", True),
    ], "01 server / 4")

    add_code_slide(p, "コード例：createStage（配信用Stage作成＋トークン）",
        ["IvsRealtimeUtils.java ─ PUBLISH+SUBSCRIBE権限・有効期限1分でStage作成"],
        [
            "ParticipantTokenConfiguration tokenConfig =",
            "    ParticipantTokenConfiguration.builder()",
            "        .capabilities(ParticipantTokenCapability.PUBLISH,",
            "                      ParticipantTokenCapability.SUBSCRIBE)",
            "        .userId(name)",
            "        .duration(PUBLISH_TOKEN_DURATION_MINUTES)   // = 1分",
            "        .build();",
            "",
            "CreateStageRequest request = CreateStageRequest.builder()",
            "        .name(name)                       // openchat-{code}-{ts}",
            "        .participantTokenConfigurations(tokenConfig)",
            "        .build();",
            "",
            "CreateStageResponse response = _client.createStage(request);",
            "return response;  // response.stage().arn() / participantTokens()",
        ],
        "01 server / 5", caption="視聴トークンは createSubscribeToken() で capabilities=SUBSCRIBE のみ・1分",
        hl={2, 3, 5, 13})

    add_table_slide(p, "① 配信開始（PerformerWsClient.onOpen）",
        ["順", "処理", "IVS"],
        [
            ["1-2", "origin(CORS)/IP取得・メンテナンス中チェック", ""],
            ["3", "token認証 → IP許可 → パフォーマ状態", ""],
            ["4-5", "通常チャット配信中でない / 既にOpenChat配信中でない", ""],
            ["6", "createRoom(performer) 内で createStage()", "●"],
            ["7", "OpenChatRoom生成・Map登録、publishToken取得", ""],
            ["8", "STAGE_READY {stageArn, publishToken, expirationTime} 返却", ""],
        ],
        "01 server / 6", col_w=[1, 8, 1],
        note="createRoomは_creatingで同時作成を排他。登録失敗時は deleteStage() でロールバックし孤児を残さない")

    add_code_slide(p, "コード例：createRoom のロールバック",
        ["OpenChatRoomManager.java ─ 作成途中で失敗したら作成済みStageを必ず削除"],
        [
            "try {",
            "    CreateStageResponse stRes = IvsRealtimeUtils.createStage(name);",
            "    stageArn     = stRes.stage().arn();",
            "    publishToken = stRes.participantTokens().get(0).token();",
            "",
            "    OpenChatRoom room = new OpenChatRoom(performer, owner, stageArn);",
            "    _openChatRoomMap.put(codeKey, room);   // ← Map登録",
            "    _creating.remove(codeKey);",
            "    return new CreateResult(room, publishToken, expirationTime);",
            "",
            "} catch (Exception e) {",
            "    _creating.remove(codeKey);",
            "    if (stageArn != null) {",
            "        IvsRealtimeUtils.deleteStage(stageArn);  // ← ロールバック",
            "    }",
            "    throw new RuntimeException(e);",
            "}",
        ],
        "01 server / 7", caption="孤児Stageを残さないための保険。createRoomは_creating.add()で同時作成を排他",
        hl={1, 13})

    add_table_slide(p, "② 視聴開始（ViewerWsClient.onOpen）",
        ["順", "処理", "IVS"],
        [
            ["1-2", "origin(CORS)チェック・performerCode必須/数値チェック", ""],
            ["3", "tokenあれば会員判定（無/無効ならゲスト）", ""],
            ["4", "getRoom(performerCode)。null なら NOT_STREAMING", ""],
            ["5", "addViewer(session, memberCode)。会員の重複視聴は拒否", ""],
            ["6", "createSubscribeToken(stageArn) で視聴トークン発行", "●"],
            ["7", "JOINED {stageArn, subscribeToken, expirationTime, isMember}", ""],
        ],
        "01 server / 8", col_w=[1, 8, 1],
        note="memberCode=0 はゲスト（重複チェック対象外）。会員は同一memberCodeで二重視聴不可")

    add_table_slide(p, "③ 視聴終了 / ④ 配信終了",
        ["場面", "処理", "IVS"],
        [
            ["③視聴終了", "onClose → removeViewer()。セッションとmemberCode解除", "なし"],
            ["③理由", "Stageは配信者のリソース。視聴者離脱では削除しない", "—"],
            ["④配信終了 1", "_openChatRoomMapから除去（新規視聴を遮断）", ""],
            ["④配信終了 2", "Owner別リストから除去", ""],
            ["④配信終了 3", "closeAllViewers() → 全視聴者へEND通知しWSクローズ", ""],
            ["④配信終了 4", "deleteStage(stageArn)", "●"],
        ],
        "01 server / 9", col_w=[2.2, 7.3, 1.2],
        note="closeRoomの順序: Map除去を最初に。closeAllViewersは_closedフラグで1度だけ実行")

    add_code_slide(p, "コード例：closeRoom（配信終了）の順序",
        ["OpenChatRoomManager.java ─ Map除去→視聴者通知→deleteStage の順"],
        [
            "// 1. まずMapから除去（以降の新規視聴を遮断）",
            "OpenChatRoom room =",
            "    (OpenChatRoom)_openChatRoomMap.remove(codeKey);",
            "if (room == null) return;        // 未登録なら何もしない",
            "",
            "// 2. Owner別リストからも除去",
            "ownerList.remove(room);",
            "",
            "// 3. 全視聴者へEND通知してWSクローズ",
            "room.closeAllViewers();",
            "",
            "// 4. 最後にIVS Stageを削除",
            "if (room.getStageArn() != null) {",
            "    IvsRealtimeUtils.deleteStage(room.getStageArn());",
            "}",
        ],
        "01 server / 10", caption="Map除去を最初にやることで teardown 中の新規視聴を防ぐ",
        hl={2, 9, 13})

    add_table_slide(p, "⑤ 孤児Stage削除（cleanupOrphans / cron毎時 0 * * * *）",
        ["順", "処理", "IVS"],
        [
            ["1", "Map上のstageArn一覧と_creatingのスナップショット取得", ""],
            ["2", "listStages() で全Stage取得", "●"],
            ["3", "openchat-接頭辞以外・Map登録済みはスキップ", ""],
            ["4", "_creating中・作成5分以内のStageはスキップ", ""],
            ["5", "残りを deleteStage() で削除", "●"],
        ],
        "01 server / 11", col_w=[1, 8, 1],
        note="サーバ再起動やクラッシュで消し損ねたStageを回収。5分以内スキップは登録前のraceでの誤削除防止")

    add_code_slide(p, "コード例：cleanupOrphans の誤削除ガード",
        ["OpenChatRoomManager.java ─ listStages()を走査し、安全なものだけdeleteStage"],
        [
            "for (StageSummary s : IvsRealtimeUtils.listStages()) {",
            "    String n = s.name();",
            "    // openchat- 接頭辞以外 / Map登録済みは対象外",
            "    if (n == null || !n.startsWith(OPENCHAT_PREFIX)) continue;",
            "    if (aliveStageArns.contains(s.arn())) continue;",
            "",
            "    // 作成中(_creating)のperformerはスキップ",
            "    Integer perfCode = parseStagePerformerCode(n);",
            "    if (creatingSnapshot.contains(perfCode)) continue;",
            "",
            "    // 作成から5分以内はスキップ（登録前のrace対策）",
            "    long ts = parseStageTimestamp(n);",
            "    if (ts > 0 && (now - ts) < SAFE_STAGE_AGE_MS) continue;",
            "",
            "    IvsRealtimeUtils.deleteStage(s.arn());   // 孤児を削除",
            "}",
        ],
        "01 server / 12", caption="cron毎時(0 * * * *)。listStages()はアカウント＋リージョン内の全Stageを返す点に注意",
        hl={0, 15})

    add_bullets_slide(p, "トークンと接続の関係 / 注意点", [
        (0, "二段構成: WSでFlax認証 → IVSトークンを払い出し（別物）", 'h'),
        (0, "パフォーマ: publishToken → PUBLISH + SUBSCRIBE", 't'),
        (0, "会員/ゲスト視聴: subscribeToken → SUBSCRIBE", 't'),
        (0, "注意点", 'h'),
        (0, "視聴終了でdeleteStageしない（消すと配信が落ちる）", 't'),
        (0, "ロールバック: createRoom途中失敗時は必ずdeleteStage", 't'),
        (0, "孤児削除のスコープ: listStages()はアカウント＋リージョン内全Stage。", 't'),
        (1, "同一アカウント/リージョンを複数環境で共有すると他環境を誤削除の恐れ。接頭辞に環境識別子を", 't'),
        (0, "ping/pong: 30秒間隔・2回未応答で約2分。検知遅延中はStage/視聴者が残る", 't'),
    ], "01 server / 13")

    save(p, "/home/tomcat/furutani/openchat-ivs/01_server_openchat_ivs.pptx")


# ============================================================
# DECK 2: Performer（配信側）
# ============================================================
def deck_performer():
    p = base()
    add_title_slide(p,
        "OpenChat × AWS IVS / 配信クライアント編",
        "openchat-performer ─ publish側SDK",
        "ローカルのカメラ/マイク映像を IVS Stage へ publish する",
        "amazon-ivs-web-broadcast v1.33.0  /  src/api/ivsStageController.js")

    add_bullets_slide(p, "概要 / import するSDK要素", [
        (0, "役割: ローカルのカメラ/マイク映像を IVS Stage へ publish（配信）", 't'),
        (0, "実装: src/api/ivsStageController.js（利用元: BroadcastComponent.vue）", 'c'),
        (0, "import するSDK要素", 'h'),
        (0, "Stage ─ Stageへの接続オブジェクト", 'c'),
        (0, "LocalStageStream ─ publishする音声/映像トラックのラッパー", 'c'),
        (0, "SubscribeType ─ 購読種別（配信側は NONE）", 'c'),
        (0, "StageEvents ─ Stageイベント定数", 'c'),
        (0, "StageConnectionState ─ 接続状態定数", 'c'),
    ], "02 performer / 2")

    add_table_slide(p, "画質プリセット（QUALITY_PRESETS）",
        ["key", "解像度", "frameRate", "maxBitrate(kbps)"],
        [
            ["low", "640 × 360", "30", "800"],
            ["mid", "854 × 480", "30", "1200"],
            ["high", "1280 × 720", "30", "2500"],
        ],
        "02 performer / 3", col_w=[1.5, 2, 2, 2.5],
        note="frameRateは最大60まで設定可能")

    add_flow_slide(p, "処理順（BroadcastComponent）", [
        ("mounted ①", "openLocalMedia(quality) ─ getUserMediaでカメラ/マイク取得→preview", False),
        ("mounted ②", "joinStage(publishToken, quality, onState) ─ Stageへ publish join", True),
        ("配信中 A", "setAudioMuted / setVideoMuted ─ 音声・映像の送出ON/OFF", False),
        ("配信中 B", "updateQuality(quality) ─ 画質変更（再joinなし）", False),
        ("beforeUnmount", "leaveStage() ─ Stage退出＋ローカルトラック停止", False),
    ], "02 performer / 4")

    add_table_slide(p, "joinStage() ─ publish専用 strategy",
        ["strategy", "値", "意味"],
        [
            ["stageStreamsToPublish", "_publishStreams", "音声/映像を配信"],
            ["shouldPublishParticipant", "true", "自身は配信者"],
            ["shouldSubscribeToParticipant", "SubscribeType.NONE", "他者は購読しない（一方向）"],
        ],
        "02 performer / 5", col_w=[3.5, 3, 4],
        note="new Stage(publishToken, strategy) → STAGE_CONNECTION_STATE_CHANGED購読 → await _stage.join()")

    add_code_slide(p, "コード例：publishトラック生成（simulcast有効化）",
        ["ivsStageController.js ─ 映像は simulcast:{enabled:true} で複数画質レイヤーを送出"],
        [
            "function _buildPublishStreams (qualityKey) {",
            "  const preset = QUALITY_PRESETS[qualityKey]",
            "  const audioTrack = new LocalStageStream(",
            "    _localStream.getAudioTracks()[0])",
            "  const videoTrack = new LocalStageStream(",
            "    _localStream.getVideoTracks()[0], {",
            "      maxBitrate:   preset.maxBitrate,",
            "      maxFramerate: preset.frameRate,",
            "      simulcast: { enabled: true }   // ← 視聴側が画質を選べる",
            "    })",
            "  _publishStreams = [audioTrack, videoTrack]  // [0]=音声 [1]=映像",
            "}",
        ],
        "02 performer / 6", caption="setAudioMuted/setVideoMuted は _publishStreams[0]/[1].setMuted() で送出ON/OFF",
        hl={8})

    add_code_slide(p, "コード例：joinStage（publish join）",
        ["ivsStageController.js ─ publish専用strategyでStageに参加し接続状態を通知"],
        [
            "export async function joinStage (publishToken, qualityKey, onState) {",
            "  _buildPublishStreams(qualityKey)",
            "",
            "  const strategy = {",
            "    stageStreamsToPublish:        () => _publishStreams,",
            "    shouldPublishParticipant:     () => true,",
            "    shouldSubscribeToParticipant: () => SubscribeType.NONE  // 一方向",
            "  }",
            "",
            "  _stage = new Stage(publishToken, strategy)",
            "  _stage.on(StageEvents.STAGE_CONNECTION_STATE_CHANGED, (s) => {",
            "    onState(s === StageConnectionState.CONNECTED ? 'connected' : ...)",
            "  })",
            "  await _stage.join()",
            "}",
        ],
        "02 performer / 7", caption="配信側は SubscribeType.NONE（他者を購読しない）。publishToken有効期限は約1分",
        hl={6, 14})

    add_bullets_slide(p, "メソッド詳細 ─ この場面でこう使う", [
        (0, "openLocalMedia(qualityKey)", 'h'),
        (0, "getUserMediaでカメラ/マイク取得（解像度/frameRateはプリセットのideal指定）。既存があればstop()して取り直す", 't'),
        (0, "joinStage() ─ _buildPublishStreams()でLocalStageStream生成", 'h'),
        (0, "音声: new LocalStageStream(audioTrack)", 'c'),
        (0, "映像: new LocalStageStream(videoTrack, { maxBitrate, maxFramerate, simulcast:{enabled:true} })", 'c'),
        (1, "simulcast有効化 → 視聴側が複数画質レイヤーから選択できる", 't'),
        (0, "updateQuality() ─ 配信中の画質変更", 'h'),
        (0, "Stageは維持。メディア取り直し→トラック差替→_stage.refreshStrategy()。再joinしないのでtoken不要", 't'),
        (0, "setAudioMuted/setVideoMuted ─ LocalStageStream.setMuted()。[0]=音声 [1]=映像", 'h'),
        (0, "leaveStage() ─ _stage.leave()＋全トラックstop()（カメラ/マイク解放）→null化", 'h'),
    ], "02 performer / 8")

    add_code_slide(p, "コード例：updateQuality（再joinなしの画質変更）",
        ["ivsStageController.js ─ Stageは維持。トラック差替→refreshStrategy()で反映"],
        [
            "/**",
            " * 配信中の画質変更。Stageは維持したまま、",
            " * トラックとbitrateを差し替える。再joinは不要なのでtoken不要。",
            " */",
            "export async function updateQuality (qualityKey) {",
            "  if (!_stage) throw new Error('not joined yet')",
            "  await _acquireMedia(qualityKey)       // メディア取り直し",
            "  _buildPublishStreams(qualityKey)      // トラック差し替え",
            "  _stage.refreshStrategy()              // ← strategyを再評価し反映",
            "  return _localStream",
            "}",
        ],
        "02 performer / 9", caption="publishTokenを使わない＝join済みStageへ再joinせず差し替えるのがポイント",
        hl={8})

    add_bullets_slide(p, "WebSocketとの関係", [
        (0, "publishTokenはWSのSTAGE_READYで受け取る", 'h'),
        (0, "{stageArn, publishToken, expirationTime}", 'c'),
        (0, "トークン有効期限は約1分 → WS接続後すみやかにjoinStageが必要", 't'),
        (0, "ログアウト/切断", 'h'),
        (0, "WS側でLogout送信 → サーバが配信終了（deleteStage）", 't'),
        (0, "クライアントはleaveStage()でローカル解放", 't'),
    ], "02 performer / 10")

    save(p, "/home/tomcat/furutani/openchat-ivs/02_performer_ivs_sdk.pptx")


# ============================================================
# DECK 3: Member（視聴側）
# ============================================================
def deck_member():
    p = base()
    add_title_slide(p,
        "OpenChat × AWS IVS / 視聴クライアント編",
        "openchat-member ─ subscribe側SDK",
        "パフォーマの映像を IVS Stage から subscribe する（自身は配信しない）",
        "amazon-ivs-web-broadcast v1.33.0  /  src/api/ivsStageController.js")

    add_bullets_slide(p, "概要 / import するSDK要素", [
        (0, "役割: パフォーマの映像を IVS Stage から subscribe（視聴）。自身は配信しない", 't'),
        (0, "実装: src/api/ivsStageController.js（利用元: ViewerComponent.vue）", 'c'),
        (0, "import するSDK要素", 'h'),
        (0, "Stage ─ Stageへの接続オブジェクト", 'c'),
        (0, "SubscribeType ─ 購読種別（視聴側は AUDIO_VIDEO）", 'c'),
        (0, "StageEvents ─ Stageイベント定数", 'c'),
        (0, "StageConnectionState ─ 接続状態定数", 'c'),
        (0, "StreamType ─ ストリーム種別（VIDEO/AUDIO判定）", 'c'),
        (0, "配信側と違い LocalStageStream は使わない（publishしないため）", 't'),
    ], "03 member / 2")

    add_flow_slide(p, "処理順（ViewerComponent）", [
        ("mounted", "joinStage(subscribeToken, onState, onStreams, onLayers) ─ subscribe-only join", True),
        ("視聴中 A", "onStreams(mediaStream) → video.srcObjectにセット（映像表示）", False),
        ("視聴中 B", "onLayers(layers) → simulcast層リスト更新（画質選択UI）", False),
        ("視聴中 C", "setPreferredLayer(label) → 視聴画質を選択", False),
        ("beforeUnmount", "leaveStage() ─ Stage退出", False),
    ], "03 member / 3")

    add_table_slide(p, "joinStage() ─ subscribe専用 strategy",
        ["strategy", "値", "意味"],
        [
            ["stageStreamsToPublish", "[]", "配信しない"],
            ["shouldPublishParticipant", "false", "視聴者"],
            ["shouldSubscribeToParticipant", "SubscribeType.AUDIO_VIDEO", "音声+映像を購読"],
            ["preferredLayerForStream", "関数", "選択中のsimulcast層を返す（未選択なら自動）"],
        ],
        "03 member / 4", col_w=[3.5, 3.2, 4],
        note="new Stage(subscribeToken, strategy) → 各イベント購読 → await _stage.join()")

    add_code_slide(p, "コード例：joinStage（subscribe-only join）",
        ["ivsStageController.js ─ 配信しないstrategy。preferredLayerForStreamで画質指定"],
        [
            "export async function joinStage (subscribeToken, onState,",
            "                                 onStreams, onLayers) {",
            "  const strategy = {",
            "    stageStreamsToPublish:        () => [],     // 配信しない",
            "    shouldPublishParticipant:     () => false,",
            "    shouldSubscribeToParticipant: () => SubscribeType.AUDIO_VIDEO,",
            "    preferredLayerForStream: (_p, stream) => {",
            "      if (!_preferredLayerLabel) return null     // 未選択=自動",
            "      if (stream.streamType !== StreamType.VIDEO) return null",
            "      return _preferredLayerLabel                // 選択中の画質層",
            "    }",
            "  }",
            "  _stage = new Stage(subscribeToken, strategy)",
            "  // ... 各イベントを購読 ...",
            "  await _stage.join()",
            "}",
        ],
        "03 member / 5", caption="視聴側は AUDIO_VIDEO を購読。LocalStageStreamは使わない（publishしないため）",
        hl={5, 6})

    add_table_slide(p, "購読イベント詳細 ─ このイベントでこう処理",
        ["イベント", "処理"],
        [
            ["STAGE_CONNECTION_STATE_CHANGED", "connecting/connected/disconnected をコールバック通知"],
            ["STAGE_PARTICIPANT_STREAMS_ADDED", "パフォーマ映像到着。MediaStream組立→onStreams。getLayers()→onLayers"],
            ["STAGE_PARTICIPANT_STREAMS_REMOVED", "配信停止。onStreams(null) / onLayers([])"],
            ["STAGE_STREAM_LAYERS_CHANGED", "simulcast層の変化。消えたら自動に戻し(refreshStrategy)→onLayers"],
        ],
        "03 member / 6", col_w=[4.5, 6.5],
        note="自身（participant.isLocal）のストリームは無視する")

    add_code_slide(p, "コード例：STREAMS_ADDED（映像受信→表示）",
        ["ivsStageController.js ─ パフォーマ映像が来たらMediaStreamを組み立てて通知"],
        [
            "_stage.on(StageEvents.STAGE_PARTICIPANT_STREAMS_ADDED,",
            "  (participant, streams) => {",
            "    if (participant && participant.isLocal) return  // 自分は無視",
            "",
            "    const mediaStream = new MediaStream()",
            "    for (let i = 0; i < streams.length; i++) {",
            "      mediaStream.addTrack(streams[i].mediaStreamTrack)",
            "    }",
            "    _onStreams(mediaStream)        // → video.srcObject にセット",
            "",
            "    const v = findVideoStream(streams)",
            "    if (v) _onLayers(v.getLayers() || [])  // simulcast層リスト",
            "  })",
        ],
        "03 member / 7", caption="STREAMS_REMOVED では onStreams(null)/onLayers([]) を通知して表示停止",
        hl={4, 8})

    add_bullets_slide(p, "STREAMS_ADDED の映像組み立て / レイヤー選択", [
        (0, "STREAMS_ADDED の映像組み立て", 'h'),
        (0, "const mediaStream = new MediaStream()", 'c'),
        (0, "for (const s of streams) mediaStream.addTrack(s.mediaStreamTrack)", 'c'),
        (0, "onStreams(mediaStream)   // video.srcObject にセット", 'c'),
        (0, "setPreferredLayer(label) ─ 視聴simulcast層（画質）を選択", 'h'),
        (0, "_preferredLayerLabel を更新し _stage.refreshStrategy() で反映", 't'),
        (0, "labelが空なら自動選択に戻る", 't'),
        (0, "preferredLayerForStreamがこの値を返すことで購読層が切り替わる", 't'),
        (0, "leaveStage() ─ _stage.leave()＋コールバック参照をnull化", 'h'),
    ], "03 member / 8")

    add_code_slide(p, "コード例：画質層の選択と自動フォールバック",
        ["ivsStageController.js ─ setPreferredLayerで選択。層が消えたら自動に戻す"],
        [
            "// 視聴画質を選択（labelが空なら自動選択に戻る）",
            "export function setPreferredLayer (label) {",
            "  _preferredLayerLabel = label || null",
            "  if (_stage) _stage.refreshStrategy()   // strategyを再評価",
            "}",
            "",
            "// 選択中の層が配信側から消えたら自動に戻す",
            "_stage.on(StageEvents.STAGE_STREAM_LAYERS_CHANGED,",
            "  (_p, stream, layers) => {",
            "    const list = layers || []",
            "    if (_preferredLayerLabel &&",
            "        !list.some(l => l.label === _preferredLayerLabel)) {",
            "      _preferredLayerLabel = null",
            "      _stage.refreshStrategy()           // 自動に戻す",
            "    }",
            "    _onLayers(list)",
            "  })",
        ],
        "03 member / 9", caption="preferredLayerForStream がこの値を返すことで購読する画質層が切り替わる",
        hl={3, 13})

    add_bullets_slide(p, "WebSocketとの関係", [
        (0, "subscribeTokenはWSのJOINEDで受け取る", 'h'),
        (0, "{stageArn, subscribeToken, expirationTime, isMember}", 'c'),
        (0, "トークン有効期限は約1分 → WS接続後すみやかにjoinStageが必要", 't'),
        (0, "配信終了・ログアウト", 'h'),
        (0, "配信終了するとサーバからOpenChatEndが届きWSクローズ → leaveStage()", 't'),
        (0, "Logout送信 → サーバがremoveViewer。視聴側ではStage削除しない（配信者のリソース）", 't'),
    ], "03 member / 10")

    save(p, "/home/tomcat/furutani/openchat-ivs/03_member_ivs_sdk.pptx")


if __name__ == "__main__":
    deck_server()
    deck_performer()
    deck_member()
