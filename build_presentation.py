# -*- coding: utf-8 -*-
"""OpenChat IVS ─ 統合プレゼン1本（導入+3章+終章）を生成。
   ・各スライドに発表者ノート（台本）を埋め込む
   ・同じ台本テキストから openchat_ivs_台本.md も出力
   対象: エンジニア / 想定15〜25分 / 1人で作ったプロジェクトの説明"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- テーマ色 ----
NAVY   = RGBColor(0x1F, 0x33, 0x55)
BLUE   = RGBColor(0x2E, 0x6F, 0xB5)
ACCENT = RGBColor(0xE8, 0x7A, 0x1E)
LIGHT  = RGBColor(0xEF, 0xF3, 0xF8)
GRAY   = RGBColor(0x55, 0x5F, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x27, 0x2E)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
FONT   = "Meiryo"

CODE_BG  = RGBColor(0x1E, 0x21, 0x2B)
CODE_TX  = RGBColor(0xE6, 0xE9, 0xEF)
CODE_CMT = RGBColor(0x7E, 0xC6, 0x99)
CODE_KEY = RGBColor(0x6A, 0xB0, 0xF3)

SW, SH = Inches(13.333), Inches(7.5)

# 台本収集用 (章, タイトル, ノート)
SCRIPT = []
_NO = [0]


def _font(run, size, color, bold=False, mono=False):
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    run.font.name = "Consolas" if mono else FONT


def _bg(s, color=WHITE):
    b = s.shapes.add_shape(1, 0, 0, SW, SH)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    b.shadow.inherit = False
    return b


def _notes(s, chapter, title, note):
    if note:
        s.notes_slide.notes_text_frame.text = note
        _NO[0] += 1
        SCRIPT.append((_NO[0], chapter, title, note))


def _header(s, title, chip=None):
    band = s.shapes.add_shape(1, 0, 0, SW, Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
    bar = s.shapes.add_shape(1, 0, Inches(1.0), SW, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(10.6), Inches(0.7)).text_frame
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tb.paragraphs[0].add_run(); r.text = title; _font(r, 25, WHITE, True)
    if chip:
        nb = s.shapes.add_textbox(Inches(10.6), Inches(0.28), Inches(2.5), Inches(0.5)).text_frame
        nb.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = nb.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = chip; _font(r, 12, RGBColor(0x9F, 0xB0, 0xC4), True)


# ---------- スライド種別 ----------

def title_slide(prs, kicker, title, subtitle, tag, note):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.25), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11), Inches(0.6)).text_frame
    r = tb.paragraphs[0].add_run(); r.text = kicker; _font(r, 18, ACCENT, True)
    tb = s.shapes.add_textbox(Inches(0.85), Inches(2.5), Inches(11.8), Inches(2.0)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = title; _font(r, 42, WHITE, True)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(4.55), Inches(11), Inches(1.2)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = subtitle; _font(r, 18, RGBColor(0xC7, 0xD3, 0xE2))
    chip = s.shapes.add_shape(1, Inches(0.9), Inches(6.2), Inches(6.0), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = BLUE; chip.line.fill.background(); chip.shadow.inherit = False
    ctf = chip.text_frame; ctf.margin_top = Pt(2); ctf.margin_bottom = Pt(2)
    r = ctf.paragraphs[0].add_run(); r.text = tag; _font(r, 13, WHITE, True)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _notes(s, "導入", title, note)
    return s


def section_slide(prs, num, title, sub, note):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.9), Inches(1.4), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11), Inches(0.8)).text_frame
    r = tb.paragraphs[0].add_run(); r.text = num; _font(r, 20, ACCENT, True)
    tb = s.shapes.add_textbox(Inches(0.85), Inches(3.1), Inches(11.8), Inches(1.2)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = title; _font(r, 38, WHITE, True)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.0)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = sub; _font(r, 16, RGBColor(0xC7, 0xD3, 0xE2))
    _notes(s, num, title, note)
    return s


def bullets_slide(prs, chapter, title, bullets, note, chip=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, title, chip)
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.9)).text_frame
    box.word_wrap = True
    first = True
    for level, text, kind in bullets:
        p = box.paragraphs[0] if first else box.add_paragraph(); first = False
        p.level = level; p.space_after = Pt(6); p.space_before = Pt(2)
        if kind == 'h':
            p.space_before = Pt(10)
            r = p.add_run(); r.text = text; _font(r, 19, BLUE, True)
        elif kind == 'c':
            r = p.add_run(); r.text = ("▶ " if level == 0 else "• ") + text; _font(r, 15, DARK, mono=True)
        else:
            r = p.add_run(); r.text = ("• " if level == 0 else "– ") + text
            _font(r, 16 if level == 0 else 15, DARK if level == 0 else GRAY)
    _notes(s, chapter, title, note)
    return s


def table_slide(prs, chapter, title, headers, rows, note, col_w=None, foot=None, chip=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, title, chip)
    nrows, ncols = len(rows) + 1, len(headers)
    width = Inches(12.13); height = Inches(0.5 * nrows)
    tbl = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.35), width, height).table
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Emu(int(int(width) * w / total))
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame; tf.word_wrap = True
        for m in ('left', 'right'): setattr(tf, 'margin_' + m, Pt(6))
        r = tf.paragraphs[0].add_run(); r.text = h; _font(r, 13, WHITE, True)
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c); cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            tf = cell.text_frame; tf.word_wrap = True
            for m in ('left', 'right'): setattr(tf, 'margin_' + m, Pt(6))
            r = tf.paragraphs[0].add_run(); r.text = val; _font(r, 12, DARK, c == 0)
    if foot:
        nb = s.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.5)).text_frame
        nb.word_wrap = True
        r = nb.paragraphs[0].add_run(); r.text = "※ " + foot; _font(r, 12, GRAY)
    _notes(s, chapter, title, note)
    return s


def flow_slide(prs, chapter, title, steps, note, chip=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, title, chip)
    n = len(steps); top = Inches(1.5); gap = Inches(0.12); avail = Inches(5.6)
    h = Emu(int((int(avail) - int(gap) * (n - 1)) / n)); y = int(top)
    for label, detail, hl in steps:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Emu(y), Inches(11.7), h)
        box.fill.solid(); box.fill.fore_color.rgb = (ACCENT if hl else LIGHT)
        box.line.color.rgb = (ACCENT if hl else RGBColor(0xCB, 0xD6, 0xE2)); box.line.width = Pt(1); box.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Pt(14)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = label; _font(r, 16, WHITE if hl else NAVY, True)
        if detail:
            r2 = p.add_run(); r2.text = "   " + detail; _font(r2, 13, WHITE if hl else GRAY)
        y += int(h) + int(gap)
    _notes(s, chapter, title, note)
    return s


def _is_comment(line):
    t = line.lstrip()
    return any(t.startswith(x) for x in ("//", "#", "*", "/*", "/**"))


def code_slide(prs, chapter, title, intro, code_lines, note, caption=None, hl=None, chip=None):
    hl = hl or set()
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, title, chip)
    top = 1.3
    if intro:
        ib = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(0.9)).text_frame
        ib.word_wrap = True; first = True
        for t in intro:
            p = ib.paragraphs[0] if first else ib.add_paragraph(); first = False
            r = p.add_run(); r.text = "• " + t; _font(r, 14, GRAY); p.space_after = Pt(2)
        top += 0.35 * len(intro) + 0.2
    box_h = 7.15 - top - (0.45 if caption else 0.1)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(top), Inches(11.93), Inches(box_h))
    panel.fill.solid(); panel.fill.fore_color.rgb = CODE_BG
    panel.line.color.rgb = RGBColor(0x3A, 0x40, 0x4E); panel.line.width = Pt(1); panel.shadow.inherit = False
    tf = panel.text_frame; tf.word_wrap = False
    tf.margin_left = Pt(14); tf.margin_right = Pt(8); tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    n = len(code_lines)
    fs = 13 if n <= 16 else (12 if n <= 20 else (11 if n <= 24 else 10))
    first = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(0); p.space_before = Pt(0); p.line_spacing = 1.0
        color = CODE_CMT if _is_comment(line) else (CODE_KEY if i in hl else CODE_TX)
        r = p.add_run(); r.text = line if line else " "
        r.font.size = Pt(fs); r.font.name = "Consolas"; r.font.color.rgb = color; r.font.bold = (i in hl)
    if caption:
        cb = s.shapes.add_textbox(Inches(0.7), Inches(top + box_h + 0.05), Inches(11.93), Inches(0.4)).text_frame
        cb.word_wrap = True
        r = cb.paragraphs[0].add_run(); r.text = caption; _font(r, 12, ACCENT, True)
    _notes(s, chapter, title, note)
    return s


def _abox(s, x, y, w, h, text, sub, fill, txt, big=18, line=None):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line; box.line.width = Pt(1.5)
    else:
        box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _font(r, big, txt, True)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = sub; _font(r, 11, txt)
    return box


def _arrow(s, x, y, w, h, color, label=None, lx=None, ly=None):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    if label:
        tb = s.shapes.add_textbox(Inches(lx), Inches(ly), Inches(w + 0.6), Inches(0.4)).text_frame
        p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label; _font(r, 12, NAVY, True)


def arch_slide(prs, title, note):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    _header(s, title)
    # Flax サーバ（上・中央）
    _abox(s, 4.9, 1.45, 3.5, 0.95, "Flax サーバ", "Stage作成 / 削除 / トークン発行", NAVY, WHITE, 18)
    # IVS Stage（中央・主役）
    _abox(s, 4.9, 3.75, 3.5, 1.5, "AWS IVS Stage", "リアルタイム配信の「場」", BLUE, WHITE, 20)
    # パフォーマ（左）
    _abox(s, 0.55, 3.95, 3.0, 1.1, "パフォーマ", "openchat-performer（publish）", ACCENT, WHITE, 17)
    # 会員/ゲスト（右）
    _abox(s, 9.75, 3.95, 3.0, 1.1, "会員 / ゲスト", "openchat-member（subscribe）", GREEN, WHITE, 17)
    # 矢印 Flax -> Stage（下向き：縦矢印は下向き図形）
    da = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.35), Inches(2.45), Inches(0.6), Inches(1.25))
    da.fill.solid(); da.fill.fore_color.rgb = RGBColor(0x9F, 0xB0, 0xC4); da.line.fill.background(); da.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(6.9), Inches(2.7), Inches(3.3), Inches(0.7)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = "createStage /\ndeleteStage"; _font(r, 11, GRAY, True)
    # 矢印 パフォーマ -> Stage（publish）
    _arrow(s, 3.6, 4.25, 1.25, 0.5, ACCENT, "publish", 3.45, 3.9)
    tb = s.shapes.add_textbox(Inches(3.35), Inches(4.75), Inches(1.7), Inches(0.35)).text_frame
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "publishToken"; _font(r, 10, GRAY, mono=True)
    # 矢印 Stage -> 会員（subscribe）
    _arrow(s, 8.45, 4.25, 1.25, 0.5, GREEN, "subscribe", 8.3, 3.9)
    tb = s.shapes.add_textbox(Inches(8.2), Inches(4.75), Inches(1.8), Inches(0.35)).text_frame
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "subscribeToken"; _font(r, 10, GRAY, mono=True)
    # 下部の WS 説明帯
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.25))
    band.fill.solid(); band.fill.fore_color.rgb = LIGHT; band.line.color.rgb = RGBColor(0xCB, 0xD6, 0xE2); band.line.width = Pt(1); band.shadow.inherit = False
    tf = band.text_frame; tf.word_wrap = True; tf.margin_left = Pt(14); tf.margin_top = Pt(8)
    r = tf.paragraphs[0].add_run(); r.text = "WebSocket（Flax）で認証 → IVSトークンを払い出す二段構成"; _font(r, 14, NAVY, True)
    for t in ["・performer WS  /ws/openChatPerformerLogin → STAGE_READY(publishToken)",
              "・viewer WS  /ws/openChatViewerLogin → JOINED(subscribeToken)　※トークン有効期限は約1分"]:
        p = tf.add_paragraph(); r = p.add_run(); r.text = t; _font(r, 12, GRAY)
    _notes(s, "導入", title, note)
    return s


# ============================================================
#  プレゼン本体
# ============================================================
def build():
    p = Presentation(); p.slide_width = SW; p.slide_height = SH

    # ---------- 0章 導入 ----------
    title_slide(p,
        "AWS IVS リアルタイム配信 / 設計と実装",
        "OpenChat ライブ配信基盤を\n1人で作った話",
        "サーバ（Java）から配信・視聴クライアント（Vue）まで、3コンポーネントを単独で設計・実装",
        "amazon-ivs-web-broadcast v1.33.0 ／ AWS SDK for Java ivsrealtime ／ ap-northeast-1",
        "本日は、OpenChatのライブ配信機能をAWS IVSで実装した話をします。特徴は、サーバのJavaから、配信クライアント・視聴クライアントのVueまで、3つのコンポーネントを全部1人で設計・実装したことです。今日はその全体像と、要所の設計判断を、実際のコードを見せながら説明します。所要は20分ほど、最後に質疑の時間を取ります。")

    bullets_slide(p, "導入", "このプロジェクトは何か / スコープ", [
        (0, "OpenChat = 1人のパフォーマが配信し、複数の会員・ゲストが視聴する一方向ライブ配信機能", 't'),
        (0, "映像配信に AWS IVS（Amazon Interactive Video Service）のリアルタイム「Stage」を採用", 't'),
        (0, "1人で担当した範囲", 'h'),
        (0, "サーバ（Flax / Java）: Stageの作成・トークン発行・削除、配信状態の管理", 't'),
        (0, "配信クライアント（openchat-performer / Vue）: カメラ映像を publish", 't'),
        (0, "視聴クライアント（openchat-member / Vue）: 映像を subscribe", 't'),
        (0, "= バックエンドのAWS連携からフロントのSDK実装まで縦に通して構築", 'c'),
    ], "まずスコープです。OpenChatは、1人のパフォーマが配信して、大勢の会員やゲストが見る、一方向のライブ配信機能です。映像の土台にはAWS IVSのリアルタイムStageを使いました。担当範囲は3つ。サーバはJavaでStageとトークンの管理。配信側と視聴側はそれぞれVueのクライアントです。ポイントは、AWSとのバックエンド連携から、フロントのSDK実装まで、縦に1本通して作った点です。")

    bullets_slide(p, "導入", "課題 ─ 何を満たす必要があったか", [
        (0, "ライブ配信の要件", 'h'),
        (0, "1配信に視聴者が多数ぶら下がる（1対多）。配信者は1人、視聴者は双方向不要", 't'),
        (0, "低遅延でリアルタイムに見せたい（録画配信ではない）", 't'),
        (0, "非機能・運用面の要件", 'h'),
        (0, "トークンが漏れても外部から勝手に使われないようにしたい", 't'),
        (0, "配信が落ちてもAWS上にゴミ（孤児リソース）を残さない", 't'),
        (0, "配信者の離脱と視聴者の離脱で、消すべきもの・消さないものを正しく分ける", 't'),
    ], "解くべき課題です。配信は1対多で、視聴者は見るだけ、双方向はいりません。そして録画ではなくリアルタイム。ここがまず技術選定に効きます。加えて運用面が重要で、トークンが漏れても悪用されない仕組み、配信が落ちてもAWSにゴミを残さない仕組み、そして配信者の離脱と視聴者の離脱で挙動を正しく分けること。これらを満たす設計が必要でした。")

    table_slide(p, "導入", "なぜ AWS IVS / Stage を選んだか",
        ["要件", "Stageの性質", "活かし方"],
        [
            ["1対多・低遅延", "リアルタイムストリーミング基盤", "WebRTCベースで低遅延に多数へ配信"],
            ["一方向で良い", "参加者ごとにpublish/subscribe権限を制御", "視聴者はSUBSCRIBEのみ付与"],
            ["回線品質の差", "simulcast（複数画質レイヤー）対応", "視聴側が端末/回線に応じ画質を選択"],
            ["短命トークン", "参加トークンの有効期限を指定可能", "1分に設定し漏洩リスクを最小化"],
            ["AWS運用", "Java/JS両SDKが提供", "サーバ＝Java、クライアント＝JS で統一"],
        ],
        "なぜIVSのStageかという話です。Stageは低遅延のリアルタイム基盤で、1対多の配信に向いています。参加者ごとにpublishとsubscribeの権限を分けられるので、一方向配信ときれいに噛み合う。simulcastで複数画質を出せるので、視聴側の回線差にも対応できる。トークンに有効期限を付けられるのでセキュリティ要件も満たせる。そしてJavaとJSの両方のSDKがあるので、サーバとクライアントを同じ概念で実装できました。要件とStageの性質が素直に対応したのが選定理由です。",
        col_w=[2.3, 4.2, 4.5])

    arch_slide(p, "アーキテクチャ全体図 ─ 3コンポーネントの関係",
        "全体像です。中央がAWS IVS Stage、配信の「場」です。左のパフォーマがそこへpublishし、右の会員・ゲストがsubscribeする。上のFlaxサーバがStageの作成・削除とトークン発行を担います。重要なのは下の二段構成で、まずWebSocketでFlaxが認証し、そのうえでIVSのトークンを払い出す。クライアントはそのトークンでStageにjoinします。トークンの有効期限は約1分と短く、これが後で出てくるセキュリティの肝です。この図の左半分が2章、右半分が3章、上のサーバが1章にあたります。")

    # ---------- 1章 サーバ ----------
    section_slide(p, "1章", "サーバ（Flax / Java）",
        "IVS処理順 ─ いつ Stage を作り、いつ消すのか。createStage / createSubscribeToken / deleteStage / listStages を呼ぶ4箇所",
        "ここから1章、サーバ側です。サーバの一番の責務は、IVSのStageをいつ作っていつ消すか、その制御です。IVSのAPIを呼ぶのは実は4箇所だけ。その4箇所と、安全に消すための工夫を見ていきます。")

    bullets_slide(p, "1章", "サーバの責務と主要コンポーネント", [
        (0, "Flaxの責務: Stageの作成・トークン発行・削除をSDK経由で管理し、配信状態をメモリ上で保持", 't'),
        (0, "主要コンポーネント", 'h'),
        (0, "IvsRealtimeUtils ─ IVS SDK呼び出しのラッパー（Stage作成/削除・トークン・一覧）", 'c'),
        (0, "OpenChatRoomManager ─ 配信のシングルトン管理（Stage作成/削除＋Room登録）", 'c'),
        (0, "OpenChatRoom ─ 1配信の状態（performer / Owner / stageArn / 視聴者）", 'c'),
        (0, "OpenChatPerformerWsClient ─ 配信用WS /ws/openChatPerformerLogin", 'c'),
        (0, "OpenChatViewerWsClient ─ 視聴用WS /ws/openChatViewerLogin", 'c'),
        (0, "対象パッケージ: jp.maru.flax.ivs / .openchat / .ws", 't'),
    ], "サーバの構成です。Flaxは、Stageの作成・トークン発行・削除をSDK越しに行い、いま誰が配信中かをメモリ上で持ちます。中心はこの2つ。IvsRealtimeUtilsがIVSのSDK呼び出しを薄くラップしたクラス。OpenChatRoomManagerが配信全体をシングルトンで管理し、Stageの作成削除とRoom登録をまとめます。1配信ぶんの状態はOpenChatRoomが持ち、配信者用と視聴者用でWebSocketのエンドポイントを分けています。",
        chip="1章 サーバ")

    table_slide(p, "1章", "この場面でこのIVS API（IvsRealtimeUtils）",
        ["メソッド", "IVS API", "用途", "トークン期限"],
        [
            ["createStage(name)", "CreateStage", "配信用Stage作成＋配信トークン取得", "publish 1分"],
            ["createSubscribeToken(arn)", "CreateParticipantToken", "視聴トークン発行", "subscribe 1分"],
            ["deleteStage(arn)", "DeleteStage", "Stage削除", "-"],
            ["listStages()", "ListStages", "孤児Stage検出用の全件取得", "-"],
        ],
        "IVSのAPIラッパーはこの4つだけです。配信開始でcreateStage、視聴開始でcreateSubscribeToken、終了でdeleteStage、そして掃除用にlistStages。注目はトークンの期限で、配信も視聴も1分。これは漏れても外から使われないための設計で、トークンはStageにjoinする初回だけ使い、join後の継続には不要だから、1分で十分なんです。Stage名は openchat-パフォーマコード-タイムスタンプ という形式にしていて、これが後の孤児削除で効いてきます。",
        col_w=[3.0, 2.4, 3.6, 1.6], chip="1章 サーバ",
        foot="トークン1分はjoin時のみ使用。Stage名 openchat-{performerCode}-{timestamp} は孤児削除でperformerCode/timestampを復元するのに使う")

    flow_slide(p, "1章", "IVS処理順サマリ ─ IVSを呼ぶのは4箇所だけ", [
        ("① 配信開始", "performer WS接続 → createStage() → publishToken返却", True),
        ("② 視聴開始", "viewer WS接続（配信中のみ）→ createSubscribeToken() → subscribeToken返却", True),
        ("③ 視聴終了", "viewer切断 → removeViewer()  ※IVS呼び出しなし", False),
        ("④ 配信終了", "performer切断 → 全視聴者へEND通知 → deleteStage()", True),
        ("⑤ 孤児削除", "cron毎時 → listStages() → 管理外Stageを deleteStage()", True),
    ], "処理順を1枚にまとめました。オレンジがIVSを呼ぶ場面です。配信開始でStageを作り、視聴開始ごとに視聴トークンを発行、配信終了でStageを削除。そしてcronで毎時、取りこぼしたStageを掃除します。グレーの③視聴終了に注目してください。ここではIVSを呼びません。Stageは配信者のリソースなので、視聴者が抜けただけで消すと配信そのものが落ちてしまう。この『消す主体は配信者だけ』という線引きが、設計上いちばん大事なポイントです。",
        chip="1章 サーバ")

    code_slide(p, "1章", "コード例：createStage（Stage作成＋トークン）",
        ["IvsRealtimeUtils.java ─ PUBLISH+SUBSCRIBE権限・有効期限1分でStageを作る"],
        [
            "ParticipantTokenConfiguration tokenConfig =",
            "    ParticipantTokenConfiguration.builder()",
            "        .capabilities(ParticipantTokenCapability.PUBLISH,",
            "                      ParticipantTokenCapability.SUBSCRIBE)",
            "        .userId(name)",
            "        .duration(PUBLISH_TOKEN_DURATION_MINUTES)   // = 1分",
            "        .build();",
            "",
            "CreateStageRequest request = CreateStageRequest.builder()",
            "        .name(name)                       // openchat-{code}-{ts}",
            "        .participantTokenConfigurations(tokenConfig)",
            "        .build();",
            "",
            "CreateStageResponse response = _client.createStage(request);",
        ],
        "createStageの実コードです。ポイントは2つ。capabilitiesでPUBLISHとSUBSCRIBEを付与しつつ、durationを1分にしている点。配信者はこの1つのトークンでStageに入ります。視聴者向けのcreateSubscribeTokenの方は、capabilitiesをSUBSCRIBEだけにして同じく1分。つまり権限の差をトークン発行の時点で付けているわけです。name にさっきの openchat- 形式の名前を渡しているのもここです。",
        caption="視聴トークンは createSubscribeToken() で capabilities=SUBSCRIBE のみ・1分",
        hl={2, 3, 5}, chip="1章 サーバ")

    code_slide(p, "1章", "コード例：createRoom のロールバック",
        ["OpenChatRoomManager.java ─ 途中で失敗したら作成済みStageを必ず消す"],
        [
            "try {",
            "    CreateStageResponse stRes = IvsRealtimeUtils.createStage(name);",
            "    stageArn     = stRes.stage().arn();",
            "    publishToken = stRes.participantTokens().get(0).token();",
            "",
            "    OpenChatRoom room = new OpenChatRoom(performer, owner, stageArn);",
            "    _openChatRoomMap.put(codeKey, room);   // ← Mapに登録",
            "    _creating.remove(codeKey);",
            "    return new CreateResult(room, publishToken, expirationTime);",
            "",
            "} catch (Exception e) {",
            "    _creating.remove(codeKey);",
            "    if (stageArn != null)",
            "        IvsRealtimeUtils.deleteStage(stageArn);   // ← ロールバック",
            "    throw new RuntimeException(e);",
            "}",
        ],
        "ここが孤児リソースを作らないための1つ目の砦です。createStageでStageを作った後、Room生成やMap登録の段階で何か失敗することがあります。そのとき、catchで作成済みのStageをdeleteStageしてから例外を投げ直す。これをやらないと、AWS上にStageだけが残ってお金もかかるしリストも汚れる。さらにcreateRoomの入口では _creating という集合にアトミックに登録していて、同じパフォーマが同時に二重で部屋を作るのを防いでいます。",
        caption="createRoom入口の _creating.add() で同一パフォーマの同時作成も排他",
        hl={1, 13}, chip="1章 サーバ")

    code_slide(p, "1章", "コード例：closeRoom（配信終了）の順序が肝",
        ["OpenChatRoomManager.java ─ Map除去 → 視聴者通知 → deleteStage の順"],
        [
            "// 1. まずMapから除去（以降の新規視聴を遮断）",
            "OpenChatRoom room =",
            "    (OpenChatRoom)_openChatRoomMap.remove(codeKey);",
            "if (room == null) return;        // 未登録なら何もしない",
            "",
            "// 2. Owner別リストからも除去",
            "ownerList.remove(room);",
            "",
            "// 3. 全視聴者へEND通知してWSクローズ",
            "room.closeAllViewers();",
            "",
            "// 4. 最後にIVS Stageを削除",
            "if (room.getStageArn() != null)",
            "    IvsRealtimeUtils.deleteStage(room.getStageArn());",
        ],
        "配信終了の処理は順番が命です。最初にMapから消す。これで終了処理中に新しい視聴者が入ってくるのを止めます。次に視聴者全員へEND通知してWSを閉じる。そして最後にStageを削除。逆順だと、Stageを消した後に滑り込んだ視聴者がエラーになったり、終了処理が二重に走ったりする。closeAllViewersは_closedフラグで一度しか動かないよう守っています。順番と冪等性をセットで設計したところです。",
        caption="Map除去を先頭に置くことで teardown 中の新規視聴を確実に止める",
        hl={2, 9, 13}, chip="1章 サーバ")

    code_slide(p, "1章", "コード例：cleanupOrphans の誤削除ガード",
        ["OpenChatRoomManager.java ─ cron毎時。安全なものだけ deleteStage"],
        [
            "for (StageSummary s : IvsRealtimeUtils.listStages()) {",
            "    String n = s.name();",
            "    if (n == null || !n.startsWith(OPENCHAT_PREFIX)) continue;",
            "    if (aliveStageArns.contains(s.arn())) continue;  // 管理中は除外",
            "",
            "    // 作成中(_creating)のパフォーマはスキップ",
            "    Integer perfCode = parseStagePerformerCode(n);",
            "    if (creatingSnapshot.contains(perfCode)) continue;",
            "",
            "    // 作成から5分以内はスキップ（登録前のrace対策）",
            "    long ts = parseStageTimestamp(n);",
            "    if (ts > 0 && (now - ts) < SAFE_STAGE_AGE_MS) continue;",
            "",
            "    IvsRealtimeUtils.deleteStage(s.arn());   // 孤児を削除",
            "}",
        ],
        "孤児を作らない2つ目の砦が、この毎時cronの掃除です。サーバが再起動したりクラッシュすると、deleteStageし損ねたStageが残る。それをlistStagesで全件見て回収します。怖いのは誤削除なので、ガードを3段かけています。管理中のものは除外、いま作成中のパフォーマのものは除外、そして作成から5分以内のものは除外。最後の5分は、createStage直後でまだMap登録前という一瞬のすきまで、現役のStageを誤って消さないための保険です。",
        caption="listStagesはアカウント＋リージョン内の全Stageを返す点に注意（終章で再登場）",
        hl={0, 14}, chip="1章 サーバ")

    # ---------- 2章 配信クライアント ----------
    section_slide(p, "2章", "配信クライアント（openchat-performer）",
        "publish側 ─ amazon-ivs-web-broadcast でローカルのカメラ/マイク映像を Stage へ送る",
        "2章は配信側のクライアントです。ここからはVueとWeb Broadcast SDKの話。カメラの映像をどうStageに送るか、画質をどう扱うかを見ます。")

    bullets_slide(p, "2章", "概要 / import するSDK要素 / 画質プリセット", [
        (0, "役割: ローカルのカメラ/マイク映像を IVS Stage へ publish（配信）", 't'),
        (0, "実装: src/api/ivsStageController.js（利用元: BroadcastComponent.vue）", 'c'),
        (0, "import するSDK要素", 'h'),
        (0, "Stage / LocalStageStream / SubscribeType / StageEvents / StageConnectionState", 'c'),
        (0, "画質プリセット（QUALITY_PRESETS）", 'h'),
        (0, "low 640×360 / mid 854×480 / high 1280×720（いずれも30fps、最大60まで可）", 't'),
        (0, "maxBitrate: low 800 / mid 1200 / high 2500 kbps", 't'),
    ], "配信クライアントの概要です。役割はシンプルで、ローカルのカメラとマイクをStageにpublishすること。SDKからはStageと、配信トラックを包むLocalStageStreamなどをimportします。画質はlow・mid・highの3プリセットを用意して、解像度とビットレートをまとめて持たせています。フレームレートは30固定ですが最大60まで設定可能です。実装はivsStageController.jsに寄せて、Vueコンポーネントからはこれを呼ぶだけにしています。",
        chip="2章 配信")

    flow_slide(p, "2章", "処理順（BroadcastComponent）", [
        ("mounted ①", "openLocalMedia(quality) ─ getUserMediaでカメラ/マイク取得 → preview表示", False),
        ("mounted ②", "joinStage(publishToken, quality, onState) ─ Stageへ publish join", True),
        ("配信中 A", "setAudioMuted / setVideoMuted ─ 音声・映像の送出ON/OFF", False),
        ("配信中 B", "updateQuality(quality) ─ 画質変更（再joinなし）", False),
        ("beforeUnmount", "leaveStage() ─ Stage退出＋ローカルトラック停止（カメラ解放）", False),
    ], "配信側のライフサイクルです。マウント時にまずカメラとマイクを取ってプレビューを出し、それからpublishTokenでStageにjoinします。配信中はミュートの切り替えと、画質の変更ができる。終了時はleaveStageでStageを抜けつつ、カメラとマイクのトラックをstopして確実に解放します。次のスライドで、このjoinと画質変更の中身を実コードで見せます。",
        chip="2章 配信")

    code_slide(p, "2章", "コード例：publishトラック生成と join",
        ["ivsStageController.js ─ simulcast有効化＋publish専用strategyでjoin"],
        [
            "// 映像トラックは simulcast を有効化して複数画質を送出",
            "const videoTrack = new LocalStageStream(",
            "  _localStream.getVideoTracks()[0], {",
            "    maxBitrate: preset.maxBitrate,",
            "    maxFramerate: preset.frameRate,",
            "    simulcast: { enabled: true }   // ← 視聴側が画質を選べる",
            "  })",
            "_publishStreams = [audioTrack, videoTrack]  // [0]=音声 [1]=映像",
            "",
            "const strategy = {",
            "  stageStreamsToPublish:        () => _publishStreams,",
            "  shouldPublishParticipant:     () => true,",
            "  shouldSubscribeToParticipant: () => SubscribeType.NONE  // 一方向",
            "}",
            "_stage = new Stage(publishToken, strategy)",
            "await _stage.join()",
        ],
        "配信の中核です。まず映像トラックを作るとき simulcast を有効にしています。これで1本の配信から複数画質のレイヤーが出て、視聴側が選べるようになる。3章への伏線です。そしてstrategyが重要で、shouldPublishParticipantはtrue、つまり自分は配信する。一方 shouldSubscribeToParticipant は SubscribeType.NONE、他人の映像は購読しない。これが一方向配信をクライアント側で表現している部分です。最後にpublishTokenでStageを作ってjoinします。",
        caption="配信側は SubscribeType.NONE。音声/映像は _publishStreams[0]/[1] として保持",
        hl={5, 13}, chip="2章 配信")

    code_slide(p, "2章", "コード例：updateQuality（再joinなしの画質変更）",
        ["ivsStageController.js ─ Stageは維持。トラック差替→refreshStrategyで反映"],
        [
            "/**",
            " * 配信中の画質変更。Stageは維持したまま、",
            " * トラックとbitrateを差し替える。再joinは不要なのでtoken不要。",
            " */",
            "export async function updateQuality (qualityKey) {",
            "  if (!_stage) throw new Error('not joined yet')",
            "  await _acquireMedia(qualityKey)       // メディアを取り直す",
            "  _buildPublishStreams(qualityKey)      // トラックを差し替える",
            "  _stage.refreshStrategy()              // ← strategyを再評価し反映",
            "  return _localStream",
            "}",
        ],
        "配信中の画質変更です。素朴にやるとStageを作り直したくなりますが、それだと一瞬切れるしトークンも要る。ここではStageは維持したまま、メディアを取り直してトラックを差し替え、refreshStrategyを呼ぶだけ。SDKが新しいトラックを拾い直してくれます。再joinしないのでpublishTokenも不要。配信を途切れさせずに画質を変えられるのがポイントです。",
        caption="join済みStageへ再joinせず差し替える＝配信を途切れさせない",
        hl={8}, chip="2章 配信")

    # ---------- 3章 視聴クライアント ----------
    section_slide(p, "3章", "視聴クライアント（openchat-member）",
        "subscribe側 ─ パフォーマの映像を Stage から受け取って表示する（自身は配信しない）",
        "最後の3章は視聴側です。配信側と対になる作りで、違いは『自分は配信しない』こと。受け取った映像の表示と、画質レイヤーの選択を見ます。")

    table_slide(p, "3章", "配信側との違い ─ strategyで対比",
        ["strategy", "配信側（performer）", "視聴側（member）"],
        [
            ["stageStreamsToPublish", "_publishStreams", "[]（配信しない）"],
            ["shouldPublishParticipant", "true", "false"],
            ["shouldSubscribeToParticipant", "SubscribeType.NONE", "SubscribeType.AUDIO_VIDEO"],
            ["preferredLayerForStream", "—", "選択中のsimulcast層を返す"],
            ["使うSDK要素", "Stage / LocalStageStream", "Stage / StreamType（LocalStageStreamなし）"],
        ],
        "まず配信側との違いを1枚で対比します。同じStageクラスを使いますが、strategyが鏡像になっています。視聴側はpublishStreamsが空、publishはfalse、そしてsubscribeはAUDIO_VIDEO、つまり音声も映像も受け取る。配信側にはなかったpreferredLayerForStreamがあって、ここでどの画質レイヤーを見るかを返します。あと視聴側はLocalStageStreamを使いません。自分は何も配信しないからです。同じ概念の裏表で、頭の切り替えが要らないのが効きました。",
        col_w=[3.4, 3.3, 4.3], chip="3章 視聴")

    code_slide(p, "3章", "コード例：joinStage（subscribe-only）と映像受信",
        ["ivsStageController.js ─ AUDIO_VIDEO購読＋STREAMS_ADDEDで映像を組み立て"],
        [
            "const strategy = {",
            "  stageStreamsToPublish:        () => [],",
            "  shouldPublishParticipant:     () => false,",
            "  shouldSubscribeToParticipant: () => SubscribeType.AUDIO_VIDEO,",
            "  preferredLayerForStream: (_p, stream) =>",
            "    _preferredLayerLabel || null     // 未選択なら自動",
            "}",
            "_stage = new Stage(subscribeToken, strategy)",
            "",
            "_stage.on(StageEvents.STAGE_PARTICIPANT_STREAMS_ADDED,",
            "  (participant, streams) => {",
            "    if (participant && participant.isLocal) return  // 自分は無視",
            "    const ms = new MediaStream()",
            "    for (const s of streams) ms.addTrack(s.mediaStreamTrack)",
            "    onStreams(ms)        // → <video>.srcObject にセット",
            "  })",
        ],
        "視聴側のjoinと映像受信です。strategyはさっきの対比どおりで、AUDIO_VIDEOを購読する設定。joinした後、パフォーマの映像が届くとSTREAMS_ADDEDイベントが飛んできます。そこで受け取ったトラックからMediaStreamを組み立てて、コールバックでVue側に渡し、videoタグのsrcObjectにセットして表示します。自分自身のストリームが混じることがあるので、isLocalなら無視しているのが地味だけど大事な一行です。",
        caption="STREAMS_REMOVED では onStreams(null) を返し、配信停止時に表示を止める",
        hl={3, 14}, chip="3章 視聴")

    code_slide(p, "3章", "コード例：画質層の選択と自動フォールバック",
        ["ivsStageController.js ─ 選択した層が消えたら自動に戻す"],
        [
            "// 視聴画質を選択（空なら自動選択へ戻る）",
            "export function setPreferredLayer (label) {",
            "  _preferredLayerLabel = label || null",
            "  if (_stage) _stage.refreshStrategy()   // strategyを再評価",
            "}",
            "",
            "// 選択中の層が配信側から消えたら自動に戻す",
            "_stage.on(StageEvents.STAGE_STREAM_LAYERS_CHANGED,",
            "  (_p, stream, layers) => {",
            "    if (_preferredLayerLabel &&",
            "        !layers.some(l => l.label === _preferredLayerLabel)) {",
            "      _preferredLayerLabel = null",
            "      _stage.refreshStrategy()           // 自動に戻す",
            "    }",
            "    onLayers(layers || [])",
            "  })",
        ],
        "配信側でsimulcastを有効にしたおかげで、視聴側は画質を選べます。setPreferredLayerでラベルを覚えてrefreshStrategyを呼ぶと、さっきのpreferredLayerForStreamがその値を返して購読する層が切り替わる。注意したのはフォールバックで、例えば配信者が画質を下げて、いま見ている層が消えることがある。そのときLAYERS_CHANGEDで検知して自動選択に戻す。これを入れないと映像が止まってしまうので、選択UIとセットで必ず要る処理です。",
        caption="配信側 simulcast → 視聴側 preferredLayer。選択UIには自動フォールバックが必須",
        hl={3, 12}, chip="3章 視聴")

    # ---------- 終章 ----------
    section_slide(p, "終章", "設計判断 / 既知の課題 / まとめ",
        "通して効いた設計判断、正直に残っている課題、そして1人で縦に作って得たこと",
        "ここまでが実装の中身でした。終章では、全体を通して効いた設計判断を整理して、正直に残っている課題も話します。")

    bullets_slide(p, "終章", "効いた設計判断（5つ）", [
        (0, "① 短命トークン（1分）", 'h'),
        (0, "WS認証→IVSトークンの二段構成。トークンはjoin時のみ使用。漏れても悪用されにくい", 't'),
        (0, "② 消す主体は配信者だけ", 'h'),
        (0, "視聴終了ではdeleteStageしない。Stageは配信者のリソースという一貫した線引き", 't'),
        (0, "③ 孤児を二重で防ぐ", 'h'),
        (0, "createRom失敗時のロールバック＋毎時cronの掃除（誤削除ガード付き）", 't'),
        (0, "④ closeRoomの順序と冪等性", 'h'),
        (0, "Map除去を先頭に。closeAllViewersは_closedフラグで1度だけ", 't'),
        (0, "⑤ simulcast × 一方向購読", 'h'),
        (0, "配信は複数画質を送り、視聴はAUDIO_VIDEOのみ購読＋画質を選択", 't'),
    ], "通して効いた判断を5つに整理しました。1つ目、短命トークン。WSで認証してから1分のIVSトークンを渡す二段構成で、漏れても悪用されにくい。2つ目、消す主体は配信者だけという線引き。これを全コードで一貫させました。3つ目、孤児リソースをロールバックと毎時掃除の二重で防ぐ。4つ目、終了処理の順序と冪等性。5つ目、配信はsimulcastで複数画質、視聴は一方向購読で画質選択。この5つが、最初に挙げた課題への答えになっています。",
        chip="終章")

    bullets_slide(p, "終章", "既知の課題 / 今後", [
        (0, "孤児削除のスコープ", 'h'),
        (0, "listStages()はAWSアカウント＋リージョン内の全Stageを返す。openchat-接頭辞に環境識別子が無い", 't'),
        (1, "同一アカウント/リージョンを複数環境で共有すると他環境の配信中Stageを誤削除する恐れ", 't'),
        (1, "対策案: 環境ごとにアカウント/リージョンを分ける、または接頭辞に環境識別子を入れる", 't'),
        (0, "切断検知の遅延", 'h'),
        (0, "ping/pongは30秒間隔・2回未応答で約2分。検知が遅れるとStageや視聴者カウントがその間残る", 't'),
        (0, "状態がメモリ常駐", 'h'),
        (0, "配信状態はシングルトンのメモリ管理。サーバ多重化やフェイルオーバ時の引き継ぎは今後の課題", 't'),
    ], "正直に残っている課題です。一番大きいのは孤児削除のスコープ。listStagesはアカウントとリージョン内の全Stageを返すのに、接頭辞に環境識別子が無いので、同じアカウントを複数環境で共有すると、他の環境の配信中Stageを孤児と誤判定して消す危険がある。対策は環境を分けるか、接頭辞に識別子を入れること。次に切断検知が最大2分遅れること、そして状態をメモリで持っているのでサーバ多重化には追加設計が要ること。ここは今後の宿題として認識しています。",
        chip="終章")

    bullets_slide(p, "終章", "まとめ", [
        (0, "OpenChatのライブ配信を、AWS IVS Stage を土台に1人で縦に構築", 't'),
        (0, "サーバ（Java）＝Stage/トークンのライフサイクル管理、孤児対策、終了処理の正しさ", 't'),
        (0, "配信（Vue）＝simulcast publish、途切れない画質変更、確実なリソース解放", 't'),
        (0, "視聴（Vue）＝一方向subscribe、映像表示、画質選択と自動フォールバック", 't'),
        (0, "3層を一貫した設計思想（短命トークン・配信者主体・孤児ゼロ）で貫けたのが成果", 'h'),
        (0, "ご清聴ありがとうございました（質疑へ）", 'c'),
    ], "まとめます。OpenChatのライブ配信を、IVS Stageを土台に、サーバから配信・視聴クライアントまで1人で縦に作りました。サーバはStageとトークンのライフサイクルと終了処理の正しさ、配信はsimulcastと途切れない画質変更、視聴は一方向購読と画質選択。そして何より、短命トークン・配信者主体・孤児ゼロという一貫した思想を3層すべてに通せたことが、1人で作ったからこその成果だと思っています。以上です、ありがとうございました。質問をどうぞ。",
        chip="終章")

    out = "/home/tomcat/furutani/openchat-ivs/openchat_ivs_presentation.pptx"
    p.slide_width = SW; p.slide_height = SH
    p.save(out)
    print("wrote", out, "(", len(p.slides._sldIdLst), "slides )")
    return p


def write_script_md():
    lines = ["# OpenChat ライブ配信基盤 ─ 発表台本（通し原稿）", "",
             "対象: エンジニア ／ 想定 15〜25分 ／ スライド `openchat_ivs_presentation.pptx` と対応", "",
             "> 各見出しの番号は本編スライドの通し番号。`[章]` はセクション。",
             "> 太字の指示（**[次へ]** 等）は口頭のつなぎの目安。", "", "---", ""]
    cur = None
    for no, chapter, title, note in SCRIPT:
        if chapter != cur:
            lines.append(f"## ── {chapter} ──"); lines.append(""); cur = chapter
        lines.append(f"### {no}. {title}")
        lines.append("")
        lines.append(note)
        lines.append("")
    md = "/home/tomcat/furutani/openchat-ivs/openchat_ivs_台本.md"
    with open(md, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print("wrote", md, "(", len(SCRIPT), "slides scripted )")


if __name__ == "__main__":
    build()
    write_script_md()
